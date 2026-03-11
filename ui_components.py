from __future__ import annotations

import base64
import hashlib
import html as html_lib
import io
import json
import os
import re
import tempfile
import textwrap
from pathlib import Path
from typing import Any

import bleach
import markdown as md
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
from bs4 import BeautifulSoup, NavigableString
from PIL import Image, ImageDraw, ImageFont, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from loaders import AssetRef, DocumentAssets
from utils import load_text_from_bytes

ALLOWED_HTML_TAGS = [
    "a",
    "abbr",
    "b",
    "blockquote",
    "br",
    "code",
    "div",
    "em",
    "h1",
    "h2",
    "h3",
    "h4",
    "h5",
    "h6",
    "hr",
    "i",
    "img",
    "li",
    "mark",
    "ol",
    "p",
    "pre",
    "span",
    "strong",
    "table",
    "tbody",
    "td",
    "th",
    "thead",
    "tr",
    "u",
    "ul",
]
ALLOWED_HTML_ATTRS = {
    "a": ["href", "title", "target"],
    "img": ["src", "alt", "title"],
    "span": ["class"],
    "div": ["class"],
    "p": ["class"],
    "code": ["class"],
    "table": ["class"],
}


def apply_theme_css(theme: str) -> None:
    if theme == "Dark":
        css = """
        <style>
        .review-card {
            border: 1px solid #3f3f46;
            border-radius: 10px;
            padding: 0.8rem;
            background: #18181b;
            color: #e4e4e7;
        }
        .review-card * { color: inherit; }
        .review-muted { color: #a1a1aa; }
        div[data-testid="stMarkdownContainer"] p,
        div[data-testid="stMarkdownContainer"] li,
        div[data-testid="stMarkdownContainer"] span,
        .stCaption,
        .stText,
        label {
            color: #e5e7eb !important;
        }
        mark { background-color: #854d0e; color: #fff7ed; padding: 0.05rem 0.2rem; border-radius: 4px; }
        </style>
        """
    elif theme == "Light":
        css = """
        <style>
        .review-card {
            border: 1px solid #d4d4d8;
            border-radius: 10px;
            padding: 0.8rem;
            background: #ffffff;
            color: #18181b;
        }
        .review-card * { color: inherit; }
        .review-muted { color: #52525b; }
        mark { background-color: #fef08a; color: #111827; padding: 0.05rem 0.2rem; border-radius: 4px; }
        </style>
        """
    else:
        css = """
        <style>
        .review-card {
            border: 1px solid color-mix(in srgb, currentColor 30%, transparent);
            border-radius: 10px;
            padding: 0.8rem;
            background: color-mix(in srgb, currentColor 5%, transparent);
        }
        .review-muted { opacity: 0.8; }
        mark { padding: 0.05rem 0.2rem; border-radius: 4px; }
        </style>
        """
    st.markdown(css, unsafe_allow_html=True)


def render_asset_status_panel(doc: DocumentAssets) -> None:
    rows = []
    for kind in ("pdf", "pptx", "html", "md", "json"):
        asset = doc.get(kind)
        rows.append(
            {
                "asset": kind.upper(),
                "status": "Loaded" if asset else "Missing",
                "name": asset.name if asset else "",
            }
        )
    st.dataframe(pd.DataFrame(rows), hide_index=True, use_container_width=True)


def render_pdf_preview(asset: AssetRef, height: int = 760, page: int = 1) -> None:
    data = asset.read_bytes()
    if not data:
        st.warning("PDF data is missing or unreadable.")
        return
    encoded = base64.b64encode(data).decode("utf-8")
    page_anchor = max(1, int(page))
    src = f"data:application/pdf;base64,{encoded}#page={page_anchor}&view=FitH"
    iframe_html = (
        f'<iframe src="{src}" '
        f'width="100%" height="{height}" type="application/pdf"></iframe>'
    )
    components.html(iframe_html, height=height + 10, scrolling=True)


def _extract_slide_text(slide: Any) -> tuple[str, list[str]]:
    title = ""
    lines: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
            continue
        try:
            text = (shape.text or "").strip()
        except Exception:
            text = ""
        if not text:
            continue
        if not title:
            title = text.splitlines()[0]
        lines.append(text)
    return title, lines


def _load_font(size: int) -> ImageFont.ImageFont:
    for font_name in ("DejaVuSans.ttf", "Arial.ttf", "LiberationSans-Regular.ttf"):
        try:
            return ImageFont.truetype(font_name, size=max(10, size))
        except OSError:
            continue
    return ImageFont.load_default()


def _parse_rgb(rgb_obj: Any) -> tuple[int, int, int] | None:
    if rgb_obj is None:
        return None
    try:
        hex_value = str(rgb_obj).strip()
        if len(hex_value) == 6:
            return (int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))
    except Exception:
        return None
    return None


def _extract_color(color_obj: Any, default: tuple[int, int, int]) -> tuple[int, int, int]:
    rgb = _parse_rgb(getattr(color_obj, "rgb", None))
    return rgb if rgb else default


def _shape_box(shape: Any, slide_w: int, slide_h: int, canvas_w: int, canvas_h: int) -> tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    x0 = max(0, int(left * canvas_w / max(1, slide_w)))
    y0 = max(0, int(top * canvas_h / max(1, slide_h)))
    x1 = max(x0 + 1, int((left + width) * canvas_w / max(1, slide_w)))
    y1 = max(y0 + 1, int((top + height) * canvas_h / max(1, slide_h)))
    return x0, y0, x1, y1


def _wrapped_lines(text: str, font: ImageFont.ImageFont, max_width_px: int, max_lines: int) -> list[str]:
    clean = text.replace("\r\n", "\n").strip()
    if not clean:
        return []
    try:
        char_px = max(6, int(font.getlength("M")))
    except Exception:
        bbox = font.getbbox("M")
        char_px = max(6, bbox[2] - bbox[0])
    width_chars = max(10, int(max_width_px / char_px))

    lines: list[str] = []
    for paragraph in clean.splitlines():
        if not paragraph.strip():
            lines.append("")
            continue
        lines.extend(textwrap.wrap(paragraph, width=width_chars) or [""])
        if len(lines) >= max_lines:
            break
    if len(lines) > max_lines:
        lines = lines[:max_lines]
    if lines:
        lines[-1] = lines[-1][: max(1, width_chars - 3)] + ("..." if len(lines[-1]) >= width_chars else "")
    return lines


def _draw_text_block(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    text: str,
    color: tuple[int, int, int] = (30, 41, 59),
) -> None:
    x0, y0, x1, y1 = box
    box_w = max(20, x1 - x0)
    box_h = max(20, y1 - y0)
    font_size = max(11, min(24, int(box_h * 0.14)))
    font = _load_font(font_size)
    lines = _wrapped_lines(text, font, max_width_px=box_w - 10, max_lines=max(1, int(box_h / (font_size + 4))))
    if not lines:
        return
    y = y0 + 4
    for line in lines:
        draw.text((x0 + 6, y), line, fill=color, font=font)
        y += font_size + 3
        if y > y1 - font_size:
            break


def _draw_picture_shape(
    canvas: Image.Image,
    shape: Any,
    box: tuple[int, int, int, int],
) -> None:
    x0, y0, x1, y1 = box
    box_w = max(1, x1 - x0)
    box_h = max(1, y1 - y0)
    try:
        blob = shape.image.blob
        picture = Image.open(io.BytesIO(blob)).convert("RGB")
    except (AttributeError, OSError, UnidentifiedImageError):
        draw = ImageDraw.Draw(canvas)
        draw.rectangle(box, fill=(230, 233, 239), outline=(148, 163, 184), width=2)
        draw.text((x0 + 8, y0 + 8), "Image", fill=(71, 85, 105), font=_load_font(14))
        return

    scale = min(box_w / picture.width, box_h / picture.height)
    target_w = max(1, int(picture.width * scale))
    target_h = max(1, int(picture.height * scale))
    resized = picture.resize((target_w, target_h))
    paste_x = x0 + (box_w - target_w) // 2
    paste_y = y0 + (box_h - target_h) // 2
    canvas.paste(resized, (paste_x, paste_y))


def _draw_table_shape(
    draw: ImageDraw.ImageDraw,
    shape: Any,
    box: tuple[int, int, int, int],
) -> None:
    x0, y0, x1, y1 = box
    draw.rectangle(box, fill=(255, 255, 255), outline=(100, 116, 139), width=2)
    row_count = 0
    col_count = 0
    try:
        table = shape.table
        row_count = len(table.rows)
        col_count = len(table.columns)
    except Exception:
        pass

    if row_count and col_count:
        row_h = max(1, int((y1 - y0) / row_count))
        col_w = max(1, int((x1 - x0) / col_count))
        for r in range(1, row_count):
            y = y0 + r * row_h
            draw.line((x0, y, x1, y), fill=(148, 163, 184), width=1)
        for c in range(1, col_count):
            x = x0 + c * col_w
            draw.line((x, y0, x, y1), fill=(148, 163, 184), width=1)
    draw.text((x0 + 6, y0 + 4), "Table", fill=(51, 65, 85), font=_load_font(12))


def _coerce_number(value: Any) -> float:
    try:
        if value is None:
            return 0.0
        return float(value)
    except (TypeError, ValueError):
        return 0.0


def _extract_chart_payload(shape: Any) -> tuple[list[str], list[tuple[str, list[float]]], str]:
    categories: list[str] = []
    series_payload: list[tuple[str, list[float]]] = []
    chart_type_name = "COLUMN"
    try:
        chart = shape.chart
        chart_type_name = str(getattr(chart, "chart_type", "")).upper()

        try:
            categories_obj = chart.plots[0].categories
            for cat in categories_obj:
                label = getattr(cat, "label", None)
                categories.append(str(label if label is not None else cat))
        except Exception:
            categories = []

        for idx, series in enumerate(chart.series):
            if idx >= 5:
                break
            name = str(getattr(series, "name", "") or f"Series {idx + 1}")
            values = [max(0.0, _coerce_number(v)) for v in list(series.values)]
            if values:
                series_payload.append((name, values))
    except Exception:
        pass

    if not categories:
        max_len = max((len(values) for _, values in series_payload), default=0)
        categories = [str(i + 1) for i in range(max_len)]

    if categories:
        target_len = len(categories)
        normalized: list[tuple[str, list[float]]] = []
        for name, values in series_payload:
            if len(values) < target_len:
                values = values + [0.0] * (target_len - len(values))
            normalized.append((name, values[:target_len]))
        series_payload = normalized

    return categories, series_payload, chart_type_name


def _draw_series_legend(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    names: list[str],
    colors: list[tuple[int, int, int]],
) -> None:
    x0, y0, x1, _ = box
    font = _load_font(12)
    x = x0 + 8
    y = y0 + 6
    for idx, name in enumerate(names[:4]):
        color = colors[idx % len(colors)]
        draw.rectangle((x, y + 3, x + 10, y + 13), fill=color, outline=(51, 65, 85), width=1)
        draw.text((x + 14, y), name[:18], fill=(30, 41, 59), font=font)
        x += 110
        if x > x1 - 110:
            break


def _draw_chart_shape(
    draw: ImageDraw.ImageDraw,
    shape: Any,
    box: tuple[int, int, int, int],
) -> None:
    x0, y0, x1, y1 = box
    draw.rectangle(box, fill=(255, 255, 255), outline=(71, 85, 105), width=2)
    categories, series_payload, chart_type_name = _extract_chart_payload(shape)
    if not series_payload:
        draw.text((x0 + 8, y0 + 8), "Chart", fill=(51, 65, 85), font=_load_font(13))
        return

    palette = [
        (37, 99, 235),
        (16, 185, 129),
        (245, 158, 11),
        (239, 68, 68),
        (139, 92, 246),
    ]

    title_box = (x0 + 6, y0 + 4, x1 - 6, y0 + 26)
    _draw_series_legend(draw, title_box, [name for name, _ in series_payload], palette)

    plot_left = x0 + 24
    plot_top = y0 + 34
    plot_right = x1 - 14
    plot_bottom = y1 - 22
    if plot_right - plot_left < 30 or plot_bottom - plot_top < 30:
        return
    draw.rectangle((plot_left, plot_top, plot_right, plot_bottom), outline=(148, 163, 184), width=1)

    value_max = max(max(values) for _, values in series_payload)
    if value_max <= 0:
        value_max = 1.0

    if "PIE" in chart_type_name:
        vals = series_payload[0][1][:8]
        total = sum(v for v in vals if v > 0)
        if total <= 0:
            draw.text((plot_left + 8, plot_top + 8), "No chart values", fill=(71, 85, 105), font=_load_font(12))
            return
        pie_size = min(plot_right - plot_left, plot_bottom - plot_top) - 10
        pie_box = (plot_left + 4, plot_top + 4, plot_left + pie_size, plot_top + pie_size)
        angle = 0.0
        for idx, value in enumerate(vals):
            if value <= 0:
                continue
            sweep = 360.0 * (value / total)
            color = palette[idx % len(palette)]
            draw.pieslice(pie_box, start=angle, end=angle + sweep, fill=color, outline=(255, 255, 255), width=1)
            angle += sweep
        return

    point_count = len(categories) if categories else len(series_payload[0][1])
    if point_count <= 0:
        return

    if "LINE" in chart_type_name:
        for s_idx, (_, values) in enumerate(series_payload):
            color = palette[s_idx % len(palette)]
            points: list[tuple[int, int]] = []
            for i, value in enumerate(values):
                x = plot_left + int(i * (plot_right - plot_left) / max(1, point_count - 1))
                y = plot_bottom - int((value / value_max) * (plot_bottom - plot_top))
                points.append((x, y))
            if len(points) >= 2:
                draw.line(points, fill=color, width=2)
            for x, y in points:
                draw.ellipse((x - 2, y - 2, x + 2, y + 2), fill=color)
    else:
        series_count = max(1, len(series_payload))
        group_w = (plot_right - plot_left) / max(1, point_count)
        bar_gap = max(1.0, group_w * 0.12)
        usable_w = max(2.0, group_w - bar_gap * 2)
        bar_w = max(2.0, usable_w / series_count)
        for i in range(point_count):
            gx = plot_left + i * group_w + bar_gap
            for s_idx, (_, values) in enumerate(series_payload):
                value = values[i] if i < len(values) else 0.0
                color = palette[s_idx % len(palette)]
                bar_h = int((value / value_max) * (plot_bottom - plot_top))
                left = int(gx + s_idx * bar_w)
                right = int(left + bar_w - 1)
                top = int(plot_bottom - bar_h)
                draw.rectangle((left, top, right, plot_bottom), fill=color, outline=(30, 41, 59), width=1)

    font = _load_font(11)
    max_labels = min(6, point_count)
    for i in range(max_labels):
        idx = int(i * (point_count - 1) / max(1, max_labels - 1))
        label = categories[idx][:10] if idx < len(categories) else str(idx + 1)
        lx = plot_left + int(idx * (plot_right - plot_left) / max(1, point_count - 1))
        draw.text((lx - 14, plot_bottom + 4), label, fill=(71, 85, 105), font=font)


def _render_slide_visual(slide: Any, slide_w: int, slide_h: int, slide_no: int, canvas_w: int = 1400) -> Image.Image:
    canvas_h = max(600, int(canvas_w * slide_h / max(1, slide_w)))
    background = (255, 255, 255)
    try:
        background = _extract_color(slide.background.fill.fore_color, default=background)
    except Exception:
        pass
    image = Image.new("RGB", (canvas_w, canvas_h), color=background)
    draw = ImageDraw.Draw(image)
    line_like_types = {
        getattr(MSO_SHAPE_TYPE, "LINE", None),
        getattr(MSO_SHAPE_TYPE, "CONNECTOR", None),
    }

    for shape in slide.shapes:
        try:
            box = _shape_box(shape, slide_w=slide_w, slide_h=slide_h, canvas_w=canvas_w, canvas_h=canvas_h)
            x0, y0, x1, y1 = box
            shape_type = getattr(shape, "shape_type", None)

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                _draw_picture_shape(image, shape, box)
                continue

            if getattr(shape, "has_table", False):
                _draw_table_shape(draw, shape, box)
                continue

            if getattr(shape, "has_chart", False):
                _draw_chart_shape(draw, shape, box)
                continue

            if shape_type in line_like_types:
                line_color = (51, 65, 85)
                try:
                    line_color = _extract_color(shape.line.fill.fore_color, default=line_color)
                except Exception:
                    pass
                draw.line((x0, y0, x1, y1), fill=line_color, width=2)
                continue

            fill_color = (242, 245, 251)
            outline_color = (148, 163, 184)
            try:
                fill_color = _extract_color(shape.fill.fore_color, default=fill_color)
            except Exception:
                pass
            try:
                outline_color = _extract_color(shape.line.fill.fore_color, default=outline_color)
            except Exception:
                pass
            draw.rectangle(box, fill=fill_color, outline=outline_color, width=2)

            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text:
                    _draw_text_block(draw, box, text)
        except Exception:
            continue

    border = (100, 116, 139)
    draw.rectangle((1, 1, canvas_w - 2, canvas_h - 2), outline=border, width=2)
    draw.text((10, 8), f"Slide {slide_no}", fill=(51, 65, 85), font=_load_font(13))
    return image


@st.cache_data(show_spinner=False)
def _render_slide_visual_cached(data: bytes, slide_index: int, canvas_w: int = 1400) -> bytes:
    presentation = Presentation(io.BytesIO(data))
    slide = presentation.slides[slide_index]
    image = _render_slide_visual(
        slide=slide,
        slide_w=int(presentation.slide_width),
        slide_h=int(presentation.slide_height),
        slide_no=slide_index + 1,
        canvas_w=canvas_w,
    )
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


@st.cache_data(show_spinner=False)
def _detect_powerpoint_com() -> tuple[bool, str]:
    if os.name != "nt":
        return False, "Native COM renderer requires Windows."
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore # noqa: F401

        pythoncom.CoInitialize()
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Quit()
        pythoncom.CoUninitialize()
        return True, ""
    except Exception as exc:
        return False, f"PowerPoint COM unavailable: {exc}"


def _sorted_slide_images(slides_dir: Path) -> list[Path]:
    candidates = [path for path in slides_dir.iterdir() if path.is_file() and path.suffix.lower() in (".png", ".jpg", ".jpeg")]

    def sort_key(path: Path) -> tuple[int, str]:
        match = re.search(r"(\d+)", path.stem)
        if match:
            return int(match.group(1)), path.name.lower()
        return 10**9, path.name.lower()

    return sorted(candidates, key=sort_key)


def _fit_export_size(
    slide_w: int,
    slide_h: int,
    max_width: int = 1920,
    max_height: int = 1920,
) -> tuple[int, int]:
    if slide_w <= 0 or slide_h <= 0:
        return max_width, int(max_width * 9 / 16)
    ratio = slide_w / slide_h
    width = max_width
    height = int(round(width / ratio))
    if height > max_height:
        height = max_height
        width = int(round(height * ratio))
    return max(320, width), max(200, height)


@st.cache_data(show_spinner=False)
def _render_pptx_via_com_cached(data: bytes, width: int = 1920, height: int = 1080) -> tuple[list[bytes], str | None]:
    if os.name != "nt":
        return [], "Native COM renderer is only available on Windows."
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except Exception as exc:
        return [], f"Missing pywin32/COM support: {exc}"

    ppt_hash = hashlib.sha256(data).hexdigest()[:16]
    with tempfile.TemporaryDirectory(prefix=f"pptx_com_{ppt_hash}_") as temp_dir:
        temp_path = Path(temp_dir)
        pptx_path = temp_path / "input.pptx"
        slides_dir = temp_path / "slides"
        slides_dir.mkdir(parents=True, exist_ok=True)
        pptx_path.write_bytes(data)

        app = None
        presentation = None
        try:
            pythoncom.CoInitialize()
            app = win32com.client.DispatchEx("PowerPoint.Application")
            # Some Office setups disallow hidden automation. Try hidden first, then visible fallback.
            try:
                app.Visible = 0
            except Exception:
                try:
                    app.Visible = 1
                except Exception:
                    pass

            try:
                presentation = app.Presentations.Open(str(pptx_path), WithWindow=False, ReadOnly=True)
            except Exception:
                # Fallback when hidden window automation is blocked by policy.
                presentation = app.Presentations.Open(str(pptx_path), WithWindow=True, ReadOnly=True)
            presentation.Export(str(slides_dir), "PNG", width, height)
        except Exception as exc:
            return [], f"PowerPoint COM export failed: {exc}"
        finally:
            try:
                if presentation is not None:
                    presentation.Close()
            except Exception:
                pass
            try:
                if app is not None:
                    app.Quit()
            except Exception:
                pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

        images: list[bytes] = []
        for image_path in _sorted_slide_images(slides_dir):
            try:
                images.append(image_path.read_bytes())
            except OSError:
                continue
        if not images:
            return [], "PowerPoint COM export produced no slide images."
        return images, None


def render_pptx_preview(asset: AssetRef) -> tuple[int, str, bool]:
    data = asset.read_bytes()
    if not data:
        st.warning("PPTX data is missing or unreadable.")
        return 1, "", False
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        st.error(f"Could not parse PPTX: {exc}")
        return 1, "", False

    slide_count = len(presentation.slides)
    if slide_count == 0:
        st.info("PPTX has no slides.")
        return 1, "", False

    com_available, com_reason = _detect_powerpoint_com()
    renderer_options = ["Python fallback"]
    if com_available:
        renderer_options.insert(0, "Native PowerPoint (COM)")

    selected_renderer = str(st.session_state.get("pptx_renderer_mode", renderer_options[0]))
    if selected_renderer not in renderer_options:
        selected_renderer = renderer_options[0]
    st.session_state["pptx_renderer_mode"] = selected_renderer

    preview_mode = str(st.session_state.get("pptx_preview_mode", "Visual fallback"))
    if preview_mode not in ("Visual fallback", "Text outline"):
        preview_mode = "Visual fallback"
        st.session_state["pptx_preview_mode"] = preview_mode
    previous_slide = int(st.session_state.get("pptx_slide_selector", 1))
    previous_slide = max(1, min(slide_count, previous_slide))
    current_slide = previous_slide

    slide_changed = False
    selected_slide = current_slide
    slide = presentation.slides[selected_slide - 1]
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    export_width, export_height = _fit_export_size(slide_w=slide_w, slide_h=slide_h, max_width=1920, max_height=1920)
    title, lines = _extract_slide_text(slide)
    if selected_renderer == "Native PowerPoint (COM)":
        com_images, com_error = _render_pptx_via_com_cached(data, width=export_width, height=export_height)
        if com_images and 0 <= selected_slide - 1 < len(com_images):
            st.image(com_images[selected_slide - 1], use_container_width=True)
        else:
            if com_error:
                st.warning(f"Native renderer unavailable: {com_error}")
            rendered_png = _render_slide_visual_cached(data, selected_slide - 1, canvas_w=1400)
            st.image(rendered_png, use_container_width=True)
    elif preview_mode == "Visual fallback":
        rendered_png = _render_slide_visual_cached(data, selected_slide - 1, canvas_w=1400)
        st.image(rendered_png, use_container_width=True)
    else:
        card_w = 1280
        card_h = max(520, int(card_w * slide_h / max(1, slide_w)))
        draw = Image.new("RGB", (card_w, card_h), color=(248, 250, 252))
        card = ImageDraw.Draw(draw)
        font = _load_font(16)
        card.rectangle((22, 22, card_w - 22, card_h - 22), outline=(30, 64, 175), width=3)
        card.text((40, 36), f"Slide {selected_slide}", fill=(30, 41, 59), font=font)
        card.text((40, 72), title or "(No title detected)", fill=(15, 23, 42), font=font)
        y = 110
        wrapped: list[str] = []
        for line in lines[:12]:
            wrapped.extend(textwrap.wrap(line, width=95) or [""])
        for line in wrapped[:24]:
            card.text((48, y), f"- {line}", fill=(51, 65, 85), font=_load_font(14))
            y += 22
            if y > card_h - 40:
                break
        st.image(draw, use_container_width=True)
    if selected_renderer == "Native PowerPoint (COM)":
        st.caption("PPTX preview uses Microsoft PowerPoint COM export (Windows only) for high-fidelity slide rendering.")
    elif com_available:
        st.caption(
            "PPTX preview uses Python fallback renderer. Switch to Native PowerPoint (COM) for highest fidelity on Windows."
        )
    else:
        st.caption(
            "PPTX preview uses Python fallback renderer (layout-aware shapes + text + embedded images). "
            "Native COM renderer is unavailable."
        )

    prev_col, label_col, next_col = st.columns([0.8, 2.4, 0.8])
    with prev_col:
        prev_clicked = st.button("<", key="pptx_prev_slide", use_container_width=True)
    with label_col:
        st.markdown(f"**Slide {current_slide} / {slide_count}**")
    with next_col:
        next_clicked = st.button(">", key="pptx_next_slide", use_container_width=True)

    target_slide = current_slide
    if prev_clicked and current_slide > 1:
        target_slide -= 1
    if next_clicked and current_slide < slide_count:
        target_slide += 1
    if target_slide != current_slide:
        st.session_state["pptx_slide_selector"] = target_slide
        slide_changed = True
        st.rerun()

    st.radio(
        "PPTX renderer",
        options=renderer_options,
        horizontal=True,
        key="pptx_renderer_mode",
        help=com_reason if (not com_available and com_reason) else "Choose COM renderer for PowerPoint-fidelity output.",
    )
    st.radio(
        "PPTX mode",
        options=["Visual fallback", "Text outline"],
        horizontal=True,
        key="pptx_preview_mode",
    )
    

    with st.expander("Extracted slide text", expanded=False):
        if lines:
            st.text("\n\n".join(lines))
        else:
            st.write("No text extracted from this slide.")

    slide_text = "\n".join([title, *lines]).strip()
    return int(selected_slide), slide_text, slide_changed


def sanitize_html(raw_html: str) -> str:
    return bleach.clean(
        raw_html,
        tags=ALLOWED_HTML_TAGS,
        attributes=ALLOWED_HTML_ATTRS,
        protocols=["http", "https", "mailto", "data"],
        strip=True,
    )


def html_to_text(raw_html: str) -> str:
    soup = BeautifulSoup(raw_html, "html.parser")
    return soup.get_text("\n")


def _focus_terms(focus_text: str) -> list[str]:
    compact = " ".join(focus_text.split())
    if not compact:
        return []
    words = [w for w in re.split(r"\s+", compact) if len(w) >= 3]
    terms: list[str] = [compact[:140], compact[:80]]
    if words:
        terms.append(" ".join(words[:8]))
        terms.append(" ".join(words[:5]))
        terms.append(words[0])
    dedup: list[str] = []
    for term in terms:
        cleaned = term.strip()
        if cleaned and cleaned.lower() not in {d.lower() for d in dedup}:
            dedup.append(cleaned)
    return dedup


def _inject_focus_mark(html_content: str, focus_text: str) -> tuple[str, bool]:
    terms = _focus_terms(focus_text)
    if not terms:
        return html_content, False
    soup = BeautifulSoup(html_content, "html.parser")
    for term in terms:
        lowered = term.lower()
        for node in soup.find_all(string=True):
            if not isinstance(node, NavigableString):
                continue
            parent_name = (node.parent.name if node.parent else "").lower()
            if parent_name in {"script", "style"}:
                continue
            raw = str(node)
            idx = raw.lower().find(lowered)
            if idx < 0:
                continue
            before = raw[:idx]
            hit = raw[idx : idx + len(term)]
            after = raw[idx + len(term) :]
            mark_tag = soup.new_tag("mark", id="sync-target")
            mark_tag.string = hit
            replacement: list[Any] = []
            if before:
                replacement.append(NavigableString(before))
            replacement.append(mark_tag)
            if after:
                replacement.append(NavigableString(after))
            node.replace_with(*replacement)
            return str(soup), True
    return html_content, False


def _render_scrollable_html(html_body: str, height: int, auto_scroll: bool = False) -> None:
    script = ""
    if auto_scroll:
        script = """
        <script>
          setTimeout(function () {
            const el = document.getElementById('sync-target');
            if (el) {
              const scroller = document.scrollingElement || document.documentElement;
              let desiredInFrameY = window.innerHeight / 2;
              try {
                const frame = window.frameElement;
                if (frame && window.parent) {
                  const rect = frame.getBoundingClientRect();
                  const parentMid = window.parent.innerHeight / 2;
                  desiredInFrameY = parentMid - rect.top;
                }
              } catch (e) {}
              const minY = 24;
              const maxY = Math.max(minY, window.innerHeight - 24);
              desiredInFrameY = Math.max(minY, Math.min(maxY, desiredInFrameY));
              const targetTop = el.offsetTop - desiredInFrameY + (el.offsetHeight / 2);
              scroller.scrollTo({ top: Math.max(0, targetTop), behavior: 'smooth' });
            }
          }, 40);
        </script>
        """
    wrapped_html = f"""
    <html>
      <head>
        <meta charset="utf-8" />
        <style>
          body {{
            margin: 0;
            padding: 0.8rem 1rem;
            background: #ffffff;
            color: #111827;
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Arial, sans-serif;
            line-height: 1.45;
          }}
          img {{ max-width: 100%; height: auto; }}
          table {{ border-collapse: collapse; width: 100%; }}
          th, td {{ border: 1px solid #d1d5db; padding: 0.35rem; }}
          pre {{
            white-space: pre-wrap;
            word-break: break-word;
            background: #f8fafc;
            border: 1px solid #e2e8f0;
            border-radius: 6px;
            padding: 0.8rem;
            font-size: 12px;
            line-height: 1.35;
          }}
          mark {{
            background: #fde68a;
            color: #111827;
            padding: 0 2px;
            border-radius: 2px;
          }}
        </style>
      </head>
      <body>{html_body}{script}</body>
    </html>
    """
    components.html(wrapped_html, height=height, scrolling=True)


def trigger_parent_center_scroll(anchor_id: str) -> None:
    html = f"""
    <script>
      setTimeout(function () {{
        try {{
          const anchor = window.parent.document.getElementById({json.dumps(anchor_id)});
          if (anchor) {{
            anchor.scrollIntoView({{ behavior: 'smooth', block: 'center' }});
          }}
        }} catch (e) {{}}
      }}, 20);
    </script>
    """
    components.html(html, height=0)


def render_html_preview(asset: AssetRef, height: int = 760, focus_text: str = "") -> str:
    raw_bytes = asset.read_bytes()
    if not raw_bytes:
        st.warning("HTML file is missing or unreadable.")
        return ""
    raw_html = load_text_from_bytes(raw_bytes)
    safe_html = sanitize_html(raw_html)
    focused_html, found = _inject_focus_mark(safe_html, focus_text=focus_text)
    _render_scrollable_html(focused_html, height=height, auto_scroll=found)
    return raw_html


def render_markdown_preview(asset: AssetRef, height: int = 760, focus_text: str = "") -> str:
    raw_bytes = asset.read_bytes()
    if not raw_bytes:
        st.warning("Markdown file is missing or unreadable.")
        return ""
    markdown_text = load_text_from_bytes(raw_bytes)
    rendered_html = md.markdown(markdown_text, extensions=["fenced_code", "tables"])
    safe_html = sanitize_html(rendered_html)
    focused_html, found = _inject_focus_mark(safe_html, focus_text=focus_text)
    _render_scrollable_html(focused_html, height=height, auto_scroll=found)
    return markdown_text


def render_json_preview(json_obj: Any, height: int = 760, focus_text: str = "") -> str:
    pretty = json.dumps(json_obj, ensure_ascii=False, indent=2)
    escaped = html_lib.escape(pretty)
    found = False
    for term in _focus_terms(focus_text):
        pattern = re.compile(re.escape(html_lib.escape(term)), flags=re.IGNORECASE)
        escaped, count = pattern.subn(r'<mark id="sync-target">\g<0></mark>', escaped, count=1)
        if count:
            found = True
            break
    _render_scrollable_html(f"<pre>{escaped}</pre>", height=height, auto_scroll=found)
    return pretty
